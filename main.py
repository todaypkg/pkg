import os
import fitz  # PyMuPDF
from telethon import TelegramClient, events
from pptx import Presentation  # للتعامل مع ملفات PPT
from docx import Document  # للتعامل مع ملفات DOCX

# جلب API ID و API Hash من متغيرات البيئة
api_id = os.environ.get("API_ID")  # جلب API ID
api_hash = os.environ.get("API_HASH")  # جلب API Hash
bot_token = os.environ.get("BOT_TOKEN")  # جلب التوكن من متغيرات البيئة

# إنشاء عميل Telethon
client = TelegramClient('bot', api_id, api_hash).start(bot_token=bot_token)

# دالة لاستخراج الصور من PDF
def extract_images_from_pdf(pdf_path: str, output_folder: str):
    """
    تستخرج الصور من ملف PDF وتخزنها في مجلد مخصص.
    :param pdf_path: مسار ملف PDF
    :param output_folder: مجلد لحفظ الصور المستخرجة
    """
    pdf_document = fitz.open(pdf_path)  # فتح ملف PDF
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)  # إنشاء مجلد الإخراج إذا لم يكن موجودًا

    # استخراج الصور من كل صفحة
    for page_number in range(len(pdf_document)):
        page = pdf_document.load_page(page_number)
        image_list = page.get_images(full=True)  # الحصول على قائمة الصور في الصفحة

        for img_index, img in enumerate(image_list):
            xref = img[0]  # المرجع الداخلي للصورة
            base_image = pdf_document.extract_image(xref)  # استخراج بيانات الصورة
            image_bytes = base_image["image"]  # بيانات الصورة (بايت)
            image_ext = base_image["ext"]  # امتداد الصورة (مثل png, jpeg)
            image_filename = f"image_page{page_number + 1}_{img_index + 1}.{image_ext}"  # اسم الملف
            image_path = os.path.join(output_folder, image_filename)  # مسار حفظ الصورة

            # حفظ الصورة في المجلد
            with open(image_path, "wb") as image_file:
                image_file.write(image_bytes)

    pdf_document.close()  # إغلاق ملف PDF

# دالة لاستخراج الصور من PPT
def extract_images_from_ppt(ppt_path: str, output_folder: str):
    """
    تستخرج الصور من ملف PPT وتخزنها في مجلد مخصص.
    :param ppt_path: مسار ملف PPT
    :param output_folder: مجلد لحفظ الصور المستخرجة
    """
    presentation = Presentation(ppt_path)  # فتح ملف PPT
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)  # إنشاء مجلد الإخراج إذا لم يكن موجودًا

    # استخراج الصور من كل شريحة
    for slide_number, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if hasattr(shape, "image"):  # إذا كان الشكل يحتوي على صورة
                image = shape.image
                image_bytes = image.blob  # بيانات الصورة (بايت)
                image_ext = image.ext  # امتداد الصورة (مثل png, jpeg)
                image_filename = f"image_slide{slide_number + 1}.{image_ext}"  # اسم الملف
                image_path = os.path.join(output_folder, image_filename)  # مسار حفظ الصورة

                # حفظ الصورة في المجلد
                with open(image_path, "wb") as image_file:
                    image_file.write(image_bytes)

# دالة لاستخراج الصور من DOCX
def extract_images_from_docx(docx_path: str, output_folder: str):
    """
    تستخرج الصور من ملف DOCX وتخزنها في مجلد مخصص.
    :param docx_path: مسار ملف DOCX
    :param output_folder: مجلد لحفظ الصور المستخرجة
    """
    document = Document(docx_path)  # فتح ملف DOCX
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)  # إنشاء مجلد الإخراج إذا لم يكن موجودًا

    # استخراج الصور من كل فقرة
    for paragraph in document.paragraphs:
        for run in paragraph.runs:
            if run.element.xpath('.//w:drawing'):  # إذا كان العنصر يحتوي على صورة
                for img in run.element.xpath('.//a:blip/@r:embed'):  # استخراج الصور
                    image_part = document.part.related_parts[img]  # جزء الصورة
                    image_bytes = image_part.blob  # بيانات الصورة (بايت)
                    image_ext = "png"  # افتراضيًا، يمكن تغييره حسب الحاجة
                    image_filename = f"image_{len(os.listdir(output_folder)) + 1}.{image_ext}"  # اسم الملف
                    image_path = os.path.join(output_folder, image_filename)  # مسار حفظ الصورة

                    # حفظ الصورة في المجلد
                    with open(image_path, "wb") as image_file:
                        image_file.write(image_bytes)

# دالة لبدء البوت
@client.on(events.NewMessage(pattern='/start'))
async def start(event):
    """
    ترحب بالمستخدم عند إرسال الأمر /start.
    """
    await event.reply("مرحباً! أرسل لي ملف (PDF, PPT, DOCX) وسأستخرج الصور منه.")

# دالة لمعالجة الملفات المرسلة
@client.on(events.NewMessage)
async def handle_file(event):
    """
    تستقبل الملف من المستخدم، تستخرج الصور، وترسلها إليه.
    """
    if event.document:
        user_id = event.sender_id  # معرف المستخدم
        file_name = event.document.attributes[0].file_name  # اسم الملف
        file_ext = os.path.splitext(file_name)[1].lower()  # امتداد الملف

        # تنزيل الملف
        file_path = await event.download_media(file=f"{user_id}_temp{file_ext}")

        output_folder = f"extracted_images_{user_id}"  # مجلد لحفظ الصور المستخرجة

        # استخراج الصور حسب نوع الملف
        if file_ext == ".pdf":
            extract_images_from_pdf(file_path, output_folder)
        elif file_ext == ".pptx":
            extract_images_from_ppt(file_path, output_folder)
        elif file_ext == ".docx":
            extract_images_from_docx(file_path, output_folder)
        else:
            await event.reply("هذا النوع من الملفات غير مدعوم.")
            return

        # إرسال الصور المستخرجة إلى المستخدم
        for image_name in os.listdir(output_folder):
            image_path = os.path.join(output_folder, image_name)
            await event.reply(file=image_path)

        # تنظيف الملفات المؤقتة
        os.remove(file_path)  # حذف الملف المؤقت
        for image_name in os.listdir(output_folder):
            os.remove(os.path.join(output_folder, image_name))  # حذف الصور
        os.rmdir(output_folder)  # حذف المجلد

# تشغيل البوت
if __name__ == "__main__":
    print("تم تشغيل البوت...")
    client.run_until_disconnected()
